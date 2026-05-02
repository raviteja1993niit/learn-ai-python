"""
builder/primitives.py
──────────────────────
Low-level drawing helpers wrapping python-pptx.
All coordinates are in inches.
"""

from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Shape type constants
_RECT   = 1   # MSO_SHAPE_TYPE RECTANGLE
_RRECT  = 5   # ROUNDED_RECTANGLE


# ── Background ────────────────────────────────────────────────────────────────

def paint_bg(slide, color: RGBColor):
    """Flood-fill a slide background with a solid colour."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


# ── Basic shapes ──────────────────────────────────────────────────────────────

def box(slide, l, t, w, h, color: RGBColor):
    """Sharp-cornered filled rectangle, no border."""
    shp = slide.shapes.add_shape(_RECT, Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp


def rbox(slide, l, t, w, h, fill: RGBColor, border: RGBColor = None, border_pt: float = 0.75):
    """Rounded rectangle with optional border."""
    shp = slide.shapes.add_shape(_RRECT, Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if border:
        shp.line.color.rgb = border
        shp.line.width = Pt(border_pt)
    else:
        shp.line.fill.background()
    return shp


def hline(slide, x, y, w, color: RGBColor, thick: float = 0.05):
    """Horizontal rule (thin filled rectangle)."""
    box(slide, x, y, w, thick, color)


def vline(slide, x, y, h, color: RGBColor, thick: float = 0.03):
    """Vertical rule."""
    box(slide, x, y, thick, h, color)


# ── Single-run text box ───────────────────────────────────────────────────────

def tb(slide, l, t, w, h, text: str,
       sz: int = 13, bold: bool = False, italic: bool = False,
       color: RGBColor = None, align=PP_ALIGN.LEFT, wrap: bool = True):
    """Add a text box with a single paragraph / single run."""
    from pptx.dml.color import RGBColor as _RGB
    if color is None:
        color = _RGB(0x00, 0x00, 0x00)
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(sz)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return txb


# ── Multi-line text box ───────────────────────────────────────────────────────
#
# rows is a list where each item is either:
#   str                      → plain text, default size/bold/color
#   (text,)                  → same
#   (text, size)
#   (text, size, bold)
#   (text, size, bold, color)

def ml(slide, l, t, w, h, rows: list,
       default_sz: int = 12, default_bold: bool = False,
       default_color: RGBColor = None, bg: RGBColor = None,
       pad: float = 0.08, wrap: bool = True):
    """Multi-line text box.  rows can mix plain strings and tuples."""
    from pptx.dml.color import RGBColor as _RGB
    if default_color is None:
        default_color = _RGB(0x00, 0x00, 0x00)

    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    if bg:
        txb.fill.solid()
        txb.fill.fore_color.rgb = bg
    tf = txb.text_frame
    tf.word_wrap = wrap
    tf.margin_left   = Inches(pad)
    tf.margin_right  = Inches(pad)
    tf.margin_top    = Inches(pad * 0.5)
    tf.margin_bottom = Inches(pad * 0.5)

    first = True
    for row in rows:
        # Normalise row to (text, sz, bold, color)
        if isinstance(row, str):
            txt, sz, bold, col = row, default_sz, default_bold, default_color
        elif len(row) == 1:
            txt, sz, bold, col = row[0], default_sz, default_bold, default_color
        elif len(row) == 2:
            txt, sz, bold, col = row[0], row[1], default_bold, default_color
        elif len(row) == 3:
            txt, sz, bold, col = row[0], row[1], row[2], default_color
        else:
            txt, sz, bold, col = row[0], row[1], row[2], row[3]

        para = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        r = para.add_run()
        r.text = txt
        r.font.size = Pt(sz)
        r.font.bold = bold
        r.font.color.rgb = col
    return txb


# ── Arrow helpers ─────────────────────────────────────────────────────────────

def arrow_right(slide, x, y, color: RGBColor, w: float = 0.38):
    return tb(slide, x, y, w, 0.36, "▶", sz=14, bold=True, color=color,
              align=PP_ALIGN.CENTER)


def arrow_down(slide, x, y, color: RGBColor):
    return tb(slide, x, y, 0.36, 0.38, "▼", sz=13, bold=True, color=color,
              align=PP_ALIGN.CENTER)


# ── Flow node (card with top accent stripe) ───────────────────────────────────

def flow_node(slide, x, y, w, h, title: str, sub: str = "",
              fill: RGBColor = None, accent: RGBColor = None,
              text_col: RGBColor = None, border: RGBColor = None):
    from builder.theme import LIGHT
    fill   = fill   or LIGHT["CARD"]
    accent = accent or LIGHT["ORANGE"]
    text_col = text_col or LIGHT["NAVY"]
    border   = border   or LIGHT["LINE"]

    rbox(slide, x, y, w, h, fill, border=border)
    hline(slide, x, y, w, accent, 0.04)
    tb(slide, x + 0.10, y + 0.08, w - 0.20, 0.38,
       title, sz=11, bold=True, color=text_col)
    if sub:
        from builder.theme import LIGHT as _L
        ml(slide, x + 0.10, y + 0.48, w - 0.20, h - 0.55,
           [(ln, 9, False, LIGHT["GRAY"]) for ln in sub.split("\n")])


# ── Standard slide chrome ─────────────────────────────────────────────────────

def slide_header(slide, T: dict, title: str, subtitle: str = ""):
    """Top bar: white strip + orange underline + title + optional subtitle."""
    box(slide, 0, 0, 13.33, 0.90, T["WHITE"])
    hline(slide, 0, 0.88, 13.33, T["ORANGE"], 0.05)
    tb(slide, 0.35, 0.08, 12.60, 0.54, title,
       sz=22, bold=True, color=T["NAVY"])
    if subtitle:
        tb(slide, 0.35, 0.60, 12.60, 0.28, subtitle,
           sz=11, color=T["GRAY"], italic=True)


def slide_footer(slide, T: dict, text: str):
    """Bottom footer bar with a thin rule and muted text."""
    hline(slide, 0, 7.32, 13.33, T["LINE"], 0.03)
    tb(slide, 0.35, 7.35, 12.60, 0.22, text, sz=8, color=T["LGRAY"])
