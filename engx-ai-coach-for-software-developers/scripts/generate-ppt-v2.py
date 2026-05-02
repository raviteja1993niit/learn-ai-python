"""
generate_ppt_v2.py  —  EngX AI Coach Final Exam  (Version 2, White Theme)

White / light professional theme.
USE CASE labels appear ONLY on: Agenda + Section Divider slides.
All other slides are clean title + content — no UC badges.

Run: python generate_ppt_v2.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Colour palette  (white theme) ─────────────────────────────────────────────
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
BG      = RGBColor(0xFA, 0xFA, 0xFD)   # slide background  (near-white)
NAVY    = RGBColor(0x0D, 0x0D, 0x35)   # primary text
DARK    = RGBColor(0x22, 0x22, 0x44)   # secondary text
GRAY    = RGBColor(0x66, 0x66, 0x88)   # muted text
LGRAY   = RGBColor(0xBB, 0xBB, 0xCC)   # very muted
ORANGE  = RGBColor(0xFA, 0x64, 0x00)   # EPAM orange — primary accent
ORANGE2 = RGBColor(0xFF, 0xE0, 0xCC)   # orange tint card
CARD    = RGBColor(0xF0, 0xF2, 0xF8)   # light card bg
CARD2   = RGBColor(0xE8, 0xF5, 0xF8)   # teal-tint card (UC2)
LINE    = RGBColor(0xDD, 0xDD, 0xEE)   # divider lines
GREEN   = RGBColor(0x2E, 0x7D, 0x32)
RED     = RGBColor(0xC6, 0x28, 0x28)
LRED    = RGBColor(0xFC, 0xEC, 0xEC)
LGREEN  = RGBColor(0xE8, 0xF5, 0xE9)
CODE_BG = RGBColor(0x1E, 0x1E, 0x38)   # keep dark for code blocks
CODE_FG = RGBColor(0xCC, 0xDD, 0xEE)
UC1_C   = RGBColor(0x15, 0x65, 0xC0)   # UC1 deep blue  (agenda + divider only)
UC2_C   = RGBColor(0x00, 0x83, 0x8F)   # UC2 teal       (agenda + divider only)
UC1_L   = RGBColor(0xE3, 0xF0, 0xFF)   # UC1 light tint
UC2_L   = RGBColor(0xE0, 0xF7, 0xF8)   # UC2 light tint

OUT = (r"C:\Users\e135408\Downloads\personal-work\learn-ai\projects"
       r"\engx-ai-coach-for-software-developers"
       r"\presentations\engx-ai-coach-final-exam-raviteja-thota-v2.pptx")

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.50)


# ─────────────────────────────────────────────────────────────────────────────
# Helpers
# ─────────────────────────────────────────────────────────────────────────────

def blank():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    return s


def box(s, l, t, w, h, fill):
    shp = s.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    return shp


def rbox(s, l, t, w, h, fill, border_col=None):
    shp = s.shapes.add_shape(5, Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if border_col:
        shp.line.color.rgb = border_col
        shp.line.width = Pt(0.75)
    else:
        shp.line.fill.background()
    return shp


def tb(s, l, t, w, h, text, sz=13, bold=False, col=NAVY,
       align=PP_ALIGN.LEFT, italic=False, wrap=True):
    t_ = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = t_.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text;  r.font.size = Pt(sz)
    r.font.bold = bold;  r.font.italic = italic
    r.font.color.rgb = col
    return t_


def ml(s, l, t, w, h, rows, dsz=12, dcol=NAVY, dbold=False, bg=None, pad=0.08):
    t_ = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    if bg:
        t_.fill.solid()
        t_.fill.fore_color.rgb = bg
    tf = t_.text_frame
    tf.word_wrap = True
    tf.margin_left   = Inches(pad)
    tf.margin_right  = Inches(pad)
    tf.margin_top    = Inches(pad * 0.5)
    tf.margin_bottom = Inches(pad * 0.5)
    first = True
    for row in rows:
        if isinstance(row, str):
            txt, sz, bd, col = row, dsz, dbold, dcol
        else:
            txt = row[0]
            sz  = row[1] if len(row) > 1 else dsz
            bd  = row[2] if len(row) > 2 else dbold
            col = row[3] if len(row) > 3 else dcol
        para = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        r = para.add_run()
        r.text = txt;  r.font.size = Pt(sz)
        r.font.bold = bd;  r.font.color.rgb = col
    return t_


def hline(s, x, y, w, col=ORANGE, thick=0.05):
    box(s, x, y, w, thick, col)


def vline(s, x, y, h, col=LINE, thick=0.03):
    box(s, x, y, thick, h, col)


def footer(s):
    hline(s, 0, 7.32, 13.33, LINE, 0.03)
    tb(s, 0.35, 7.35, 9.0, 0.22,
       "EPAM Proprietary & Confidential.  |  EngX AI Coach — Final Exam  |  May 2026",
       sz=8, col=LGRAY)


def slide_header(s, title, subtitle=""):
    """Clean white-theme slide header — no UC badges."""
    box(s, 0, 0, 13.33, 0.90, WHITE)
    hline(s, 0, 0.88, 13.33, ORANGE, 0.05)
    tb(s, 0.35, 0.08, 12.60, 0.54, title, sz=22, bold=True, col=NAVY)
    if subtitle:
        tb(s, 0.35, 0.60, 12.60, 0.28, subtitle, sz=11, col=GRAY, italic=True)


def flow_node(s, x, y, w, h, title, sub="", fill=CARD, accent=ORANGE, text_col=NAVY):
    rbox(s, x, y, w, h, fill, border_col=LINE)
    hline(s, x, y, w, accent, 0.04)
    tb(s, x+0.10, y+0.08, w-0.20, 0.38, title, sz=11, bold=True, col=text_col)
    if sub:
        ml(s, x+0.10, y+0.48, w-0.20, h-0.55,
           [(ln, 9, False, GRAY) for ln in sub.split("\n")])


def arrow_r(s, x, y, w=0.38):
    tb(s, x, y, w, 0.36, "▶", sz=14, bold=True, col=ORANGE, align=PP_ALIGN.CENTER)


def arrow_d(s, x, y):
    tb(s, x, y, 0.36, 0.38, "▼", sz=13, bold=True, col=ORANGE, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — COVER
# ─────────────────────────────────────────────────────────────────────────────
s = blank()
box(s, 0, 0, 0.22, 7.50, ORANGE)          # left accent stripe
box(s, 0.22, 0, 13.11, 7.50, WHITE)       # white body

# Title block
box(s, 0.55, 0.70, 8.0, 0.06, ORANGE)
tb(s, 0.55, 0.92, 11.5, 0.72,
   "How to Develop Your Own", sz=30, bold=True, col=NAVY)
tb(s, 0.55, 1.70, 11.5, 0.72,
   "Multi-Agentic Solution", sz=30, bold=True, col=ORANGE)
box(s, 0.55, 2.55, 8.0, 0.06, LINE)

tb(s, 0.55, 2.78, 10.0, 0.42,
   "Raviteja Thota  |  Senior Software Engineer  |  EPAM", sz=14, col=NAVY)
tb(s, 0.55, 3.24, 10.0, 0.36,
   "EngX AI Coach — Final Exam Presentation  |  May 2026", sz=12, col=GRAY)

# Two use-case pills
box(s, 0.55, 3.90, 3.20, 0.52, UC1_L)
hline(s, 0.55, 3.90, 3.20, UC1_C, 0.05)
tb(s, 0.70, 3.98, 3.00, 0.36,
   "Use Case 1  ·  Multi-Agentic Solution", sz=11, bold=True, col=UC1_C)

box(s, 0.55, 4.56, 3.20, 0.52, UC2_L)
hline(s, 0.55, 4.56, 3.20, UC2_C, 0.05)
tb(s, 0.70, 4.64, 3.00, 0.36,
   "Use Case 2  ·  RAG Evolution Journey", sz=11, bold=True, col=UC2_C)

# Right decorative area
box(s, 9.20, 0.60, 3.90, 6.10, CARD)
tb(s, 9.45, 1.10, 3.50, 0.55,
   "EngX AI Coach", sz=16, bold=True, col=ORANGE, align=PP_ALIGN.CENTER)
hline(s, 9.45, 1.72, 3.30, ORANGE, 0.04)
ml(s, 9.45, 1.88, 3.45, 4.50, [
    ("Final Exam", 13, True, NAVY),
    ("", 6, False, NAVY),
    ("Topic:", 11, True, GRAY),
    ("Multi-Agentic", 12, True, NAVY),
    ("Solution Design", 12, True, NAVY),
    ("", 6, False, NAVY),
    ("2 Use Cases", 11, True, GRAY),
    ("1. Multi-Agent System", 11, False, NAVY),
    ("2. RAG Evolution", 11, False, NAVY),
    ("", 6, False, NAVY),
    ("18 slides", 11, False, GRAY),
    ("20-25 min presentation", 11, False, GRAY),
], dsz=11, dcol=NAVY)

tb(s, 0.55, 7.18, 9.0, 0.25,
   "EPAM Proprietary & Confidential.", sz=9, col=LGRAY)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — AGENDA  (UC1 / UC2 labels HERE only)
# ─────────────────────────────────────────────────────────────────────────────
s = blank()
slide_header(s, "Agenda")
footer(s)

# UC1 column
box(s, 0.35, 1.05, 5.95, 5.90, UC1_L)
hline(s, 0.35, 1.05, 5.95, UC1_C, 0.06)
box(s, 0.35, 1.05, 2.50, 0.50, UC1_C)
tb(s, 0.37, 1.07, 2.46, 0.46,
   "Use Case 1", sz=12, bold=True, col=WHITE, align=PP_ALIGN.CENTER)
tb(s, 3.00, 1.10, 3.20, 0.38,
   "Multi-Agentic Solution", sz=12, bold=True, col=UC1_C)

uc1_items = [
    "01   What is a Multi-Agent System?",
    "02   Architecture Flow Diagram",
    "03   Agent Roles & Responsibilities",
    "04   Code: Agent Definition",
    "05   Real-World Results",
]
for i, item in enumerate(uc1_items):
    tb(s, 0.55, 1.70 + i * 1.00, 5.60, 0.80, item, sz=13, col=UC1_C if i==0 else NAVY, bold=(i==0))

# Divider
vline(s, 6.55, 1.05, 5.90, LGRAY, 0.04)

# UC2 column
box(s, 6.98, 1.05, 5.95, 5.90, UC2_L)
hline(s, 6.98, 1.05, 5.95, UC2_C, 0.06)
box(s, 6.98, 1.05, 2.50, 0.50, UC2_C)
tb(s, 7.00, 1.07, 2.46, 0.46,
   "Use Case 2", sz=12, bold=True, col=WHITE, align=PP_ALIGN.CENTER)
tb(s, 9.62, 1.10, 3.20, 0.38,
   "RAG Evolution Journey", sz=12, bold=True, col=UC2_C)

uc2_items = [
    "06   Why RAG Needed to Evolve",
    "07   RAG Evolution Flow  (V1 › V2 › V3)",
    "08   Agentic RAG Architecture",
    "09   Code Demo: ask.py",
    "10   Comparison & When to Use Each",
]
for i, item in enumerate(uc2_items):
    tb(s, 7.18, 1.70 + i * 1.00, 5.60, 0.80, item, sz=13, col=UC2_C if i==0 else NAVY, bold=(i==0))

tb(s, 0.35, 7.00, 12.60, 0.28,
   "Practical Task  ·  Key Takeaways  ·  Q&A", sz=10, col=GRAY)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — ABOUT ME
# ─────────────────────────────────────────────────────────────────────────────
s = blank()
slide_header(s, "About Me")
footer(s)

ml(s, 0.35, 1.05, 6.0, 6.00, [
    ("Raviteja Thota", 22, True, NAVY),
    ("Senior Software Engineer  @  EPAM", 13, False, DARK),
    ("Hyderabad, India", 11, False, GRAY),
    ("", 5, False, NAVY),
    ("Stack:", 12, True, ORANGE),
    ("Angular  ·  Java / Spring Boot  ·  AWS  ·  Microservices", 11, False, DARK),
    ("", 5, False, NAVY),
    ("Certifications:", 12, True, ORANGE),
    ("  Claude Certified Architect – Foundations  (Anthropic 2026)", 11, False, DARK),
    ("  Microsoft AZ-400 DevOps Engineer Expert  (2026)", 11, False, DARK),
    ("  AWS Certified Developer – Associate", 11, False, DARK),
    ("  Google Professional Cloud Developer", 11, False, DARK),
    ("  Oracle OCI Generative AI Professional", 11, False, DARK),
    ("  AWS AI Practitioner", 11, False, DARK),
    ("", 5, False, NAVY),
    ("India Innovation Enabler  (Feb 2026)", 12, True, ORANGE),
    ("  Smart Mapper + ATF→Flow Framework at LEAP (global)", 11, False, GRAY),
])
vline(s, 6.58, 1.05, 6.00, LINE, 0.04)
ml(s, 6.80, 1.05, 6.18, 6.00, [
    ("Apps Built for This Exam:", 13, True, NAVY),
    ("", 5, False, NAVY),
    ("agentic-ai-toggle-management", 12, True, DARK),
    ("  Orchestrator + 5 agents, Confluence MCP", 11, False, GRAY),
    ("  Java toggle analysis & code refactoring", 11, False, GRAY),
    ("", 4, False, NAVY),
    ("vector-rag-app-v1", 12, True, DARK),
    ("  ChromaDB + sentence-transformers + BM25", 11, False, GRAY),
    ("", 4, False, NAVY),
    ("vector-less-rag-app-v2", 12, True, DARK),
    ("  BM25 + FastAPI, no vector DB", 11, False, GRAY),
    ("", 4, False, NAVY),
    ("vector-less-rag-app", 12, True, DARK),
    ("  BM25 + Streamlit, multi-query search", 11, False, GRAY),
    ("", 4, False, NAVY),
    ("vector-less-rag-agentic", 12, True, DARK),
    ("  LLM as retriever + .github/agents", 11, False, GRAY),
    ("", 4, False, NAVY),
    ("transcript-notes-generator", 12, True, DARK),
    ("  faster-whisper + Copilot LLM + agents", 11, False, GRAY),
])

# ─────────────────────────────────────────────────────────────────────────────
# USE CASE 1 — SECTION DIVIDER  (UC label ONLY here + agenda)
# ─────────────────────────────────────────────────────────────────────────────
s = blank()
box(s, 0, 0, 13.33, 7.50, UC1_C)          # full blue background
box(s, 0.28, 0, 13.05, 7.50, WHITE)       # white panel offset
box(s, 0.28, 0, 0.10, 7.50, UC1_C)       # blue left bar on white
hline(s, 0.50, 2.90, 12.60, UC1_C, 0.06)

box(s, 0.55, 1.00, 2.80, 0.62, UC1_C)
tb(s, 0.60, 1.04, 2.70, 0.54,
   "Use Case 1", sz=16, bold=True, col=WHITE, align=PP_ALIGN.CENTER)

tb(s, 0.55, 1.82, 12.20, 1.10,
   "Multi-Agentic Solution", sz=38, bold=True, col=NAVY)
tb(s, 0.55, 2.98, 12.20, 0.45,
   "Real enterprise implementation — MPGS Toggle Management System", sz=14, col=GRAY, italic=True)
hline(s, 0.55, 3.56, 4.0, UC1_C, 0.04)
ml(s, 0.55, 3.78, 12.0, 2.80, [
    ("What you will learn:", 13, True, NAVY),
    ("  ›  What makes a system truly multi-agentic (vs single LLM)", 12, False, DARK),
    ("  ›  Orchestrator + Sub-agent architecture and flow", 12, False, DARK),
    ("  ›  How to define agents as code  (.github/agents)", 12, False, DARK),
    ("  ›  Results from a live production implementation", 12, False, DARK),
])
tb(s, 0.55, 7.18, 9.0, 0.25,
   "EPAM Proprietary & Confidential.", sz=9, col=LGRAY)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — SINGLE LLM vs MULTI-AGENT
# ─────────────────────────────────────────────────────────────────────────────
s = blank()
slide_header(s, "What is a Multi-Agent System?", "And why a single LLM is not enough")
footer(s)

# Left: Single LLM
box(s, 0.35, 1.05, 5.90, 5.95, LRED)
hline(s, 0.35, 1.05, 5.90, RED, 0.05)
tb(s, 0.52, 1.14, 5.55, 0.44, "Single LLM  —  Limitations", sz=14, bold=True, col=RED)
ml(s, 0.52, 1.66, 5.55, 5.10, [
    ("One model handles everything.", 12, True, NAVY),
    ("", 5, False, NAVY),
    ("Context window overflow", 12, True, RED),
    ("  16,000+ Java files cannot fit in one prompt.", 11, False, DARK),
    ("", 5, False, NAVY),
    ("No specialisation", 12, True, RED),
    ("  A generalist LLM performs worse than", 11, False, DARK),
    ("  a focused specialist agent.", 11, False, DARK),
    ("", 5, False, NAVY),
    ("Sequential only — no parallelism", 12, True, RED),
    ("  One task at a time, slow for complex work.", 11, False, DARK),
    ("", 5, False, NAVY),
    ("Single point of failure", 12, True, RED),
    ("  One prompt failure = entire pipeline fails.", 11, False, DARK),
    ("", 5, False, NAVY),
    ("Hard to maintain", 12, True, RED),
    ("  Monolithic prompts break as scope grows.", 11, False, DARK),
])

# VS badge
rbox(s, 6.04, 3.76, 1.22, 0.62, WHITE, border_col=ORANGE)
tb(s, 6.06, 3.80, 1.18, 0.54, "VS", sz=20, bold=True, col=ORANGE, align=PP_ALIGN.CENTER)

# Right: Multi-Agent
box(s, 7.40, 1.05, 5.58, 5.95, LGREEN)
hline(s, 7.40, 1.05, 5.58, GREEN, 0.05)
tb(s, 7.58, 1.14, 5.22, 0.44, "Multi-Agent System  —  Advantages", sz=14, bold=True, col=GREEN)
ml(s, 7.58, 1.66, 5.22, 5.10, [
    ("Orchestrator routes. Agents specialise.", 12, True, NAVY),
    ("", 5, False, NAVY),
    ("Separation of concerns", 12, True, GREEN),
    ("  Each agent owns one domain.", 11, False, DARK),
    ("", 5, False, NAVY),
    ("Composable & extensible", 12, True, GREEN),
    ("  Add or swap agents without redesign.", 11, False, DARK),
    ("", 5, False, NAVY),
    ("Resumable long-running workflows", 12, True, GREEN),
    ("  Shared state persists across turns.", 11, False, DARK),
    ("", 5, False, NAVY),
    ("MCP tool composability", 12, True, GREEN),
    ("  Each agent gets its own tool set.", 11, False, DARK),
    ("", 5, False, NAVY),
    ("Traceable & testable", 12, True, GREEN),
    ("  Every handoff is logged. Test each agent.", 11, False, DARK),
])

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — ARCHITECTURE FLOW DIAGRAM
# ─────────────────────────────────────────────────────────────────────────────
s = blank()
slide_header(s, "Architecture Flow: Orchestrator + 5 Sub-Agents",
             "MPGS Agentic Toggle Management System")
footer(s)

# User node
NW, NH = 2.18, 0.60
flow_node(s, 5.57, 1.05, NW, NH, "👤  User Command",
          "Plain English or /slash-cmd", fill=CARD, accent=ORANGE)
arrow_d(s, 6.57, 1.67)

# Orchestrator node
OW, OH = 5.60, 1.10
flow_node(s, 3.85, 2.10, OW, OH,
          "🎯  TOGGLE-ORCHESTRATOR  (Router)",
          "Routes commands · Answers simple queries inline · NEVER does leaf work · Enforces guardrails",
          fill=UC1_L, accent=UC1_C, text_col=UC1_C)

# Connector down + horizontal bar
box(s, 6.60, 3.22, 0.04, 0.30, UC1_C)
box(s, 1.52, 3.52, 10.32, 0.04, UC1_C)

# 5 sub-agent columns
agents = [
    ("SA-1\nANALYSIS ENGINE",    "Java source scan\nON/OFF impact\nS2A enrichment"),
    ("SA-2\nREPORT WRITER",      "6 report types\nMarkdown output\nCoverage check"),
    ("SA-3\nCONFLUENCE ENQUIRY", "4-tier MCP search\nPage metadata\nBatch fetch"),
    ("SA-4\nADVISOR\n★ on-request","Decomp planning\nJava refactoring\nToggle removal"),
    ("SA-5\nWORKSPACE MANAGER",  "Session state\nRegistry mgmt\nRepo onboarding"),
]
SA_W, SA_H = 2.04, 2.70
SA_Y = 3.92
for i, (name, desc) in enumerate(agents):
    sx = 0.40 + i * 2.58
    box(s, sx + SA_W/2 - 0.02, 3.52, 0.04, 0.38, UC1_C)
    arrow_d(s, sx + SA_W/2 - 0.18, 3.84)
    fill_ = UC1_L if i != 3 else ORANGE2
    ac_   = UC1_C if i != 3 else ORANGE
    flow_node(s, sx, SA_Y, SA_W, SA_H, name, desc, fill=fill_, accent=ac_,
              text_col=UC1_C if i != 3 else ORANGE)

# Tools bar
box(s, 0.40, 6.68, 12.52, 0.42, CARD)
hline(s, 0.40, 6.68, 12.52, LINE, 0.03)
tb(s, 0.56, 6.74, 12.20, 0.32,
   "Shared Tools:  read_file  ·  grep_search  ·  run_subagent  ·  "
   "insert_edit_into_file  ·  mcp-confluence/getPage  ·  run_in_terminal",
   sz=9, col=GRAY)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7 — AGENT ROLES TABLE
# ─────────────────────────────────────────────────────────────────────────────
s = blank()
slide_header(s, "Agent Roles & Responsibilities",
             "Each agent owns exactly one domain — clear input/output contracts")
footer(s)

hcols = ["Agent", "Role", "Key Tools", "Output"]
cxs   = [0.35, 2.65, 6.60, 10.10]
cws   = [2.26, 3.90, 3.46, 3.10]
box(s, 0.35, 1.05, 12.62, 0.44, NAVY)
for x, w, h in zip(cxs, cws, hcols):
    tb(s, x+0.10, 1.08, w-0.12, 0.38, h, sz=11, bold=True, col=WHITE)

rows_a = [
    ("toggle-orchestrator", "Session entry point. Routes all commands. NEVER does leaf work.",
     "run_subagent, read_file", "Delegated or inline answer"),
    ("SA-1  analysis-engine", "Java source scan across 7 repos. Builds ON/OFF business & field-level impact.",
     "grep_search, read_file, run_in_terminal", "Analysis CSV per toggle"),
    ("SA-2  report-writer", "Generates 6 structured report types from analysis data.",
     "read_file, insert_edit_into_file", "Markdown report files"),
    ("SA-3  confluence-enquiry", "4-tier Confluence search. Extracts page metadata for each toggle.",
     "mcp-confluence/getPage, searchContent", "Results CSV + metadata"),
    ("SA-4  advisor  ★", "Plans & executes Java toggle decomposition (on-request only).",
     "read_file, apply_patch, grep_search", "Decomp plan / refactored Java"),
    ("SA-5  workspace-manager", "Session state, toggle registry, workspace cleanup, repo onboarding.",
     "read_file, create_file, insert_edit", "Registry + session docs"),
]
for ri, row in enumerate(rows_a):
    y = 1.53 + ri * 0.96
    bg = WHITE if ri % 2 == 0 else CARD
    box(s, 0.35, y, 12.62, 0.92, bg)
    hline(s, 0.35, y + 0.88, 12.62, LINE, 0.03)
    tc0 = ORANGE if ri == 0 else DARK
    for xi, (x, w, cell) in enumerate(zip(cxs, cws, row)):
        tb(s, x+0.10, y+0.10, w-0.14, 0.72, cell, sz=9,
           col=tc0 if xi == 0 else DARK, bold=(xi == 0))

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8 — CODE: AGENT DEFINITION
# ─────────────────────────────────────────────────────────────────────────────
s = blank()
slide_header(s, "Code: Defining an Agent",
             ".github/agents/toggle-orchestrator.agent.md")
footer(s)

code1 = """\
---
name: toggle-orchestrator
description: Session entry point and router for the toggle management system.
tools:
  - read_file
  - run_subagent          # delegates to specialist sub-agents
  - grep_search
  - mcp-confluence/getPage
  - insert_edit_into_file
---

# Toggle Orchestrator

## Routing Strategy

### Tier 1 — Handle Inline  (no sub-agent delegation)
| Request                       | Orchestrator Action                         |
|-------------------------------|---------------------------------------------|
| Single toggle data query      | Read analysis CSV, reply as Markdown        |
| Analysis status check         | Count registry rows, answer inline          |
| Toggle already analysed       | Display existing rows from analysis CSV     |

### Tier 2 — Delegate to Sub-Agent
| User Command                           | Sub-Agent Invoked              |
|----------------------------------------|--------------------------------|
| analyze toggle <name>                  | toggle-analysis-engine  SA-1   |
| generate report                        | toggle-report-writer    SA-2   |
| fetch confluence data for <name>       | toggle-confluence-enquiry SA-3 |
| plan toggle decomposition for <repo>   | toggle-advisor          SA-4   |
| manage session / onboard repo          | toggle-workspace-manager SA-5  |
"""
box(s, 0.35, 1.05, 12.62, 6.10, CODE_BG)
tb(s, 0.55, 1.14, 12.22, 5.94, code1, sz=10, col=CODE_FG)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 9 — RESULTS
# ─────────────────────────────────────────────────────────────────────────────
s = blank()
slide_header(s, "Results: What the System Delivered",
             "MPGS Payment Gateway — production implementation")
footer(s)

metrics = [
    ("100+", "Feature toggles\nanalysed"),
    ("7",    "Java repos\nscanned"),
    ("16K+", "Java files\nprocessed"),
    ("6",    "Report types\nauto-generated"),
    ("5",    "Specialist\nagents"),
]
MW = 2.36
for i, (num, label) in enumerate(metrics):
    mx = 0.35 + i * 2.60
    box(s, mx, 1.05, MW, 1.90, CARD)
    hline(s, mx, 1.05, MW, ORANGE, 0.06)
    tb(s, mx, 1.18, MW, 0.90, num, sz=34, bold=True, col=ORANGE, align=PP_ALIGN.CENTER)
    tb(s, mx+0.10, 2.02, MW-0.20, 0.80, label, sz=10, col=DARK, align=PP_ALIGN.CENTER)

# Before / After flow
box(s, 0.35, 3.12, 5.88, 3.62, LRED)
hline(s, 0.35, 3.12, 5.88, RED, 0.05)
tb(s, 0.52, 3.20, 5.52, 0.42, "Before  (Manual)", sz=13, bold=True, col=RED)
ml(s, 0.52, 3.70, 5.52, 2.70, [
    ("Engineer reads Java files manually", 11, False, DARK),
    ("  ▼", 10, False, GRAY),
    ("Checks Confluence (often stale)", 11, False, DARK),
    ("  ▼", 10, False, GRAY),
    ("Writes impact notes by hand", 11, False, DARK),
    ("  ▼", 10, False, GRAY),
    ("Days per toggle  ·  Error-prone", 12, True, RED),
])

tb(s, 6.32, 4.65, 0.78, 0.50, "▶", sz=26, bold=True, col=ORANGE, align=PP_ALIGN.CENTER)

box(s, 7.22, 3.12, 5.75, 3.62, LGREEN)
hline(s, 7.22, 3.12, 5.75, GREEN, 0.05)
tb(s, 7.38, 3.20, 5.42, 0.42, "After  (Multi-Agent System)", sz=13, bold=True, col=GREEN)
ml(s, 7.38, 3.70, 5.42, 2.70, [
    ('User types:  "analyze toggle <name>"', 11, True, DARK),
    ("  ▼", 10, False, GRAY),
    ("Orchestrator routes to SA-1 (analysis)", 11, False, DARK),
    ("  ▼", 10, False, GRAY),
    ("SA-1 scans repos · SA-3 fetches docs", 11, False, DARK),
    ("  ▼", 10, False, GRAY),
    ("Minutes per toggle  ·  Reproducible", 12, True, GREEN),
])

# ─────────────────────────────────────────────────────────────────────────────
# USE CASE 2 — SECTION DIVIDER  (UC label ONLY here + agenda)
# ─────────────────────────────────────────────────────────────────────────────
s = blank()
box(s, 0, 0, 13.33, 7.50, UC2_C)
box(s, 0.28, 0, 13.05, 7.50, WHITE)
box(s, 0.28, 0, 0.10, 7.50, UC2_C)
hline(s, 0.50, 2.90, 12.60, UC2_C, 0.06)

box(s, 0.55, 1.00, 2.80, 0.62, UC2_C)
tb(s, 0.60, 1.04, 2.70, 0.54,
   "Use Case 2", sz=16, bold=True, col=WHITE, align=PP_ALIGN.CENTER)

tb(s, 0.55, 1.82, 12.20, 1.10,
   "RAG Evolution Journey", sz=38, bold=True, col=NAVY)
tb(s, 0.55, 2.98, 12.20, 0.45,
   "Vector RAG  ›  Vector-less RAG  ›  Agentic RAG", sz=14, col=GRAY, italic=True)
hline(s, 0.55, 3.56, 4.0, UC2_C, 0.04)
ml(s, 0.55, 3.78, 12.0, 2.80, [
    ("What you will learn:", 13, True, NAVY),
    ("  ›  Why RAG architectures evolved across 3 generations", 12, False, DARK),
    ("  ›  Trade-offs: embeddings vs BM25 vs LLM-as-retriever", 12, False, DARK),
    ("  ›  How to build Agentic RAG with .github/agents", 12, False, DARK),
    ("  ›  Which approach to use for your project", 12, False, DARK),
])
tb(s, 0.55, 7.18, 9.0, 0.25,
   "EPAM Proprietary & Confidential.", sz=9, col=LGRAY)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 11 — WHY RAG EVOLVED
# ─────────────────────────────────────────────────────────────────────────────
s = blank()
slide_header(s, "Why RAG Needed to Evolve",
             "Each generation solved the main pain of the previous one")
footer(s)

pains = [
    ("Pain 1", "Vector DB Overhead",
     "ChromaDB setup, embedding model download,\n"
     "infra cost, slow first run, version-lock.\n"
     "Barrier: too complex for quick proofs-of-concept."),
    ("Pain 2", "Semantic Gap in BM25",
     "BM25 finds keyword matches but misses meaning.\n"
     "Query 'payment failure' misses 'transaction declined'.\n"
     "Barrier: poor recall on business documents."),
    ("Pain 3", "Retriever Code Complexity",
     "Chunkers, retrievers, BM25 indexes — all custom code.\n"
     "Adding a new doc format = rewriting chunker + parser.\n"
     "Barrier: hard to maintain, hard to extend."),
]
for i, (label, title, body) in enumerate(pains):
    px = 0.35 + i * 4.33
    box(s, px, 1.05, 4.10, 3.02, LRED)
    hline(s, px, 1.05, 4.10, RED, 0.05)
    tb(s, px+0.14, 1.14, 3.84, 0.32, label, sz=10, bold=True, col=RED)
    tb(s, px+0.14, 1.50, 3.84, 0.42, title, sz=13, bold=True, col=NAVY)
    tb(s, px+0.14, 1.98, 3.84, 1.92, body, sz=10, col=DARK)

tb(s, 0.35, 4.22, 12.62, 0.38,
   "How each generation solved these pains:", sz=13, bold=True, col=NAVY)

sols = [
    ("V1  ›  V2", "Dropped ChromaDB + embeddings entirely.\nBM25 on raw pages — zero infra cost."),
    ("V2  ›  V3", "Dropped BM25 + all retriever code.\nLLM reads full JSON doc — zero retrieval code."),
    ("V3 + Agents", "Added .github/agents for multi-turn Q&A.\nOrchestrate: convert → ask → summarise."),
]
for i, (label, body) in enumerate(sols):
    sx = 0.35 + i * 4.33
    box(s, sx, 4.68, 4.10, 2.40, CARD)
    hline(s, sx, 4.68, 4.10, ORANGE, 0.05)
    tb(s, sx+0.14, 4.78, 3.84, 0.40, label, sz=12, bold=True, col=ORANGE)
    tb(s, sx+0.14, 5.25, 3.84, 1.68, body, sz=11, col=DARK)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 12 — RAG EVOLUTION FLOW DIAGRAM
# ─────────────────────────────────────────────────────────────────────────────
s = blank()
slide_header(s, "RAG Evolution Flow: V1  ›  V2  ›  V3",
             "Each version removes one layer of complexity — same output, less code")
footer(s)

def rag_node(s, x, y, w, h, title, sub, removed=False):
    fill = LRED if removed else CARD
    ac   = RED  if removed else ORANGE
    rbox(s, x, y, w, h, fill, border_col=LINE)
    hline(s, x, y, w, ac, 0.04)
    tb(s, x+0.10, y+0.08, w-0.20, 0.36, title, sz=10, bold=True,
       col=RED if removed else DARK)
    if sub:
        tb(s, x+0.10, y+0.46, w-0.20, h-0.55, sub, sz=9, col=GRAY, wrap=True)

NW2, NH2 = 3.86, 0.72
VS = [0.35, 4.74, 9.13]
col_titles = ["V1 — Vector RAG", "V2 — Vector-less RAG", "V3 — Agentic RAG"]
col_apps   = ["vector-rag-app-v1", "vector-less-rag-app-v2", "vector-less-rag-agentic"]
for ci, (cx, ct, ca) in enumerate(zip(VS, col_titles, col_apps)):
    col_acc = ORANGE if ci < 2 else UC2_C
    tb(s, cx, 1.05, NW2, 0.34, ct, sz=11, bold=True, col=col_acc)
    tb(s, cx, 1.40, NW2, 0.25, ca, sz=9, col=GRAY, italic=True)

# Row labels
row_labels = ["Document Upload", "Chunker", "Embedder", "Vector DB / Index", "Retrieve", "LLM Answer"]
row_ys     = [1.70, 2.48, 3.26, 4.04, 4.82, 5.60]

# V1 column — all present
v1_subs = ["PDF / DOCX / XLSX", "512-token chunks", "sentence-transformers", "ChromaDB", "Cosine sim top-K", "GitHub Copilot LLM"]
for rl, ry, sub in zip(row_labels, row_ys, v1_subs):
    rag_node(s, 0.35, ry, NW2, 0.72, rl, sub, removed=False)
    if ry < 5.60:
        arrow_d(s, 0.35 + NW2/2 - 0.18, ry + 0.72)

# V2 column — Chunker & Embedder & ChromaDB REMOVED
v2_data = [
    ("Document Upload",    "PDF / DOCX / XLSX", False),
    ("Chunker",            "REMOVED", True),
    ("Embedder",           "REMOVED", True),
    ("BM25 Index",         "Keyword index over pages", False),
    ("Retrieve top-K",     "Multi-query BM25", False),
    ("LLM Answer",         "GitHub Copilot LLM", False),
]
for (title, sub, rem), ry in zip(v2_data, row_ys):
    rag_node(s, 4.74, ry, NW2, 0.72, title, sub, removed=rem)
    if ry < 5.60 and not rem:
        arrow_d(s, 4.74 + NW2/2 - 0.18, ry + 0.72)

# V3 column — Chunker, Embedder, ChromaDB, BM25 REMOVED; add JSON Tree
v3_data = [
    ("Document Upload",    "PDF / DOCX / XLSX", False),
    ("convert.py",         "Extracts JSON tree structure", False),
    ("Embedder",           "REMOVED", True),
    ("Retriever / Index",  "REMOVED", True),
    ("LLM reads full doc", "Full JSON context, no retrieval", False),
    ("LLM Answer",         "Section-aware response", False),
]
for (title, sub, rem), ry in zip(v3_data, row_ys):
    rag_node(s, 9.13, ry, NW2, 0.72, title, sub, removed=rem)
    if ry < 5.60 and not rem:
        arrow_d(s, 9.13 + NW2/2 - 0.18, ry + 0.72)

# Arrows between columns
for ay in [3.75, 3.75]:
    pass
tb(s, 4.24, 3.85, 0.46, 0.46, "›", sz=22, bold=True, col=ORANGE, align=PP_ALIGN.CENTER)
tb(s, 8.63, 3.85, 0.46, 0.46, "›", sz=22, bold=True, col=ORANGE, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 13 — AGENTIC RAG ARCHITECTURE FLOW
# ─────────────────────────────────────────────────────────────────────────────
s = blank()
slide_header(s, "Agentic RAG: Full Architecture Flow",
             "vector-less-rag-agentic  —  .github/agents orchestrate the pipeline")
footer(s)

# Top row: 7 pipeline nodes
pipeline = [
    ("👤 User",        "Uploads doc\nasks question"),
    ("📄 Document",    "PDF / DOCX\n/ XLSX"),
    ("🔄 convert.py",  "JSON tree\nextraction"),
    ("📋 JSON Tree",   "sections & tables\nstructured"),
    ("🤖 ask.py",      "Sends full JSON\n+ question to LLM"),
    ("🧠 Copilot LLM", "Reads & retrieves\nfrom full context"),
    ("💬 Answer",      "Section-aware\nresponse"),
]
PNW, PNH = 1.72, 1.50
PY = 1.10
for i, (title, sub) in enumerate(pipeline):
    px = 0.32 + i * 1.86
    fill = UC2_L if i in (4, 5) else CARD
    ac   = UC2_C if i in (4, 5) else ORANGE
    tc   = UC2_C if i in (4, 5) else DARK
    flow_node(s, px, PY, PNW, PNH, title, sub, fill=fill, accent=ac, text_col=tc)
    if i < len(pipeline) - 1:
        arrow_r(s, px + PNW + 0.02, PY + PNH/2 - 0.18, 0.38)

# Agents layer header
tb(s, 0.32, 2.80, 12.68, 0.36,
   "GitHub Copilot  .github/agents  —  Orchestration Layer", sz=12, bold=True, col=NAVY)
hline(s, 0.32, 3.18, 12.68, ORANGE, 0.04)

# 4 agent nodes
ag_nodes = [
    ("doc-processor.md\n+ pdf-processor.md", "Runs convert.py\nRoutes to qa-agent after"),
    ("qa-agent.md",            "Calls ask.py\nMulti-turn Q&A\nSession memory"),
    ("summarizer.md",          "Calls summarize.py\nSection-by-section\nsummaries"),
    ("copilot-instructions.md","Shared context\nModel selection\nAuth setup"),
]
AGW, AGH = 2.96, 2.72
for i, (name, desc) in enumerate(ag_nodes):
    ax = 0.32 + i * 3.22
    fill = UC2_L if i < 3 else CARD
    ac   = UC2_C if i < 3 else ORANGE
    tc   = UC2_C if i < 3 else DARK
    flow_node(s, ax, 3.32, AGW, AGH, name, desc, fill=fill, accent=ac, text_col=tc)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 14 — CODE DEMO: ask.py
# ─────────────────────────────────────────────────────────────────────────────
s = blank()
slide_header(s, "Code Demo: ask.py — The LLM IS the Retriever",
             "vector-less-rag-agentic  |  No BM25  ·  No embeddings  ·  Only 3 Python files")
footer(s)

code2 = "\n".join([
    "# ask.py  --  LLM acts as the retrieval agent  (vector-less-rag-agentic)",
    "# Traditional RAG: embed -> ChromaDB -> BM25 -> retrieve chunks -> LLM",
    "# Agentic RAG:     parse doc to JSON tree -> send FULL JSON -> LLM retrieves",
    "",
    'def ask(json_path: str, question: str, model: str = "gpt-4o-mini") -> str:',
    '    doc     = json.loads(Path(json_path).read_text(encoding="utf-8"))',
    "    context = build_context(doc)        # full JSON  -->  readable Markdown sections",
    "",
    "    client = openai.OpenAI(",
    '        base_url = "https://models.inference.ai.azure.com",',
    "        api_key  = get_token_or_raise(),   # reads from 'gh auth token' automatically",
    "    )",
    "    response = client.chat.completions.create(",
    "        model    = model,",
    "        messages = [",
    '            {"role": "system", "content": SYSTEM_PROMPT},',
    '            {"role": "user",   "content": f"{context}\\n\\nQuestion: {question}"},',
    "        ],",
    "        temperature = 0.2,",
    "        max_tokens  = 1024,",
    "    )",
    "    return response.choices[0].message.content",
    "",
    "",
    "def build_context(doc: dict) -> str:",
    "    # Convert JSON tree to readable Markdown -- no chunking, no BM25.",
    "    parts = []",
    "    for sec in doc.get('sections', []):",
    "        parts.append(f\"## {sec['heading']}\\n{sec['content']}\")",
    "    return '\\n\\n'.join(parts)",
    "",
    "",
    "# .github/agents/qa-agent.md orchestrates the session:",
    "#   User question  ->  agent  ->  run ask.py  ->  return answer  ->  multi-turn",
])
box(s, 0.35, 1.05, 12.62, 6.10, CODE_BG)
tb(s, 0.52, 1.14, 12.22, 5.92, code2, sz=10, col=CODE_FG)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 15 — COMPARISON TABLE
# ─────────────────────────────────────────────────────────────────────────────
s = blank()
slide_header(s, "Comparison: V1  vs  V2  vs  V3  —  When to Use Which?")
footer(s)

cols15  = ["", "V1  —  Vector RAG", "V2  —  Vector-less RAG", "V3  —  Agentic RAG"]
c15xs   = [0.35, 3.18, 6.40, 9.65]
c15ws   = [2.79, 3.18, 3.21, 3.50]
box(s, 0.35, 1.05, 12.62, 0.44, NAVY)
for x, w, h in zip(c15xs, c15ws, cols15):
    tb(s, x+0.10, 1.08, w-0.12, 0.38, h, sz=11, bold=True, col=WHITE)

rows15 = [
    ("Setup complexity",          "High — ChromaDB, GPU, model", "Low — BM25 only",      "Very low — 3 files"),
    ("Infrastructure required",   "Vector DB + GPU",             "None",                 "None"),
    ("Semantic search quality",   "Best (embeddings)",           "Partial (keyword)",    "Best (LLM reasoning)"),
    ("Large docs  (>100 pages)",  "Best — chunked",              "Good",                 "Limited by context"),
    ("Small–medium docs",         "Overkill",                    "Good",                 "Best"),
    ("Code to maintain",          "Chunker+Embedder+DB",         "Retriever+Index",      "Zero retrieval code"),
    ("Multi-turn agent support",  "Manual wiring",               "Manual wiring",        "Built-in via agents"),
    ("Best for",                  "Enterprise, large corpus",    "Mid-size, fast POC",   "Structured docs, RAG POC"),
]
for ri, row in enumerate(rows15):
    y = 1.53 + ri * 0.70
    bg = WHITE if ri % 2 == 0 else CARD
    box(s, 0.35, y, 12.62, 0.66, bg)
    hline(s, 0.35, y+0.62, 12.62, LINE, 0.03)
    for xi, (x, w, cell) in enumerate(zip(c15xs, c15ws, row)):
        col = GRAY if xi == 0 else \
              GREEN if ("None" in cell or "Zero" in cell or "Best" in cell or "Built" in cell) else \
              RED   if ("Overkill" in cell or "Manual" in cell or "Limited" in cell) else DARK
        tb(s, x+0.10, y+0.08, w-0.14, 0.52, cell, sz=9,
           col=col, bold=(xi == 0))

tb(s, 0.35, 7.06, 12.62, 0.24,
   "Rule: Can you fit the full doc in context? → V3 Agentic.  Large corpus? → V1.  Quick prototype? → V2.",
   sz=10, bold=False, col=GRAY)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 16 — PRACTICAL TASK
# ─────────────────────────────────────────────────────────────────────────────
s = blank()
slide_header(s, "Practical Task  —  Build Your Own  (15 minutes)")
footer(s)

# Left task
box(s, 0.35, 1.05, 5.95, 5.90, CARD)
hline(s, 0.35, 1.05, 5.95, ORANGE, 0.06)
tb(s, 0.52, 1.18, 5.62, 0.44,
   "Task A  —  Multi-Agent Code Review", sz=13, bold=True, col=NAVY)
ml(s, 0.52, 1.70, 5.62, 5.00, [
    ("Build a 2-agent code review system:", 12, True, DARK),
    ("", 5, False, DARK),
    ("Step 1 — Orchestrator agent", 11, True, ORANGE),
    ("  .github/agents/code-review-orchestrator.agent.md", 10, False, GRAY),
    ("  Accept file path. Delegate. Return summary.", 10, False, DARK),
    ("  Tools: run_subagent", 10, False, GRAY),
    ("", 5, False, DARK),
    ("Step 2 — Analyser agent", 11, True, ORANGE),
    ("  .github/agents/code-analyser.agent.md", 10, False, GRAY),
    ("  Find: functions>30 lines, no error handling.", 10, False, DARK),
    ("  Tools: read_file, grep_search", 10, False, GRAY),
    ("  Output: Markdown table — issue | line | fix", 10, False, GRAY),
    ("", 5, False, DARK),
    ("Step 3 — Test it", 11, True, ORANGE),
    ("  @code-review-orchestrator  analyse src/app.py", 10, False, GRAY),
    ("", 5, False, DARK),
    ("Bonus: 3rd agent creates JIRA tickets", 11, True, NAVY),
])

vline(s, 6.55, 1.05, 5.90, LINE, 0.04)

# Right task
box(s, 6.98, 1.05, 5.95, 5.90, CARD)
hline(s, 6.98, 1.05, 5.95, ORANGE, 0.06)
tb(s, 7.14, 1.18, 5.62, 0.44,
   "Task B  —  Agentic RAG on Your Doc", sz=13, bold=True, col=NAVY)
ml(s, 7.14, 1.70, 5.62, 5.00, [
    ("Build an agentic RAG for your document:", 12, True, DARK),
    ("", 5, False, DARK),
    ("Step 1 — Convert your document", 11, True, ORANGE),
    ("  python convert.py  myreport.pdf", 10, False, GRAY),
    ("  Open the JSON output and inspect it.", 10, False, DARK),
    ("", 5, False, DARK),
    ("Step 2 — Ask via CLI", 11, True, ORANGE),
    ('  python ask.py  myreport.json', 10, False, GRAY),
    ('    "What are the key findings?"', 10, False, DARK),
    ("  Try at least 3 different questions.", 10, False, GRAY),
    ("", 5, False, DARK),
    ("Step 3 — Create qa-agent.md", 11, True, ORANGE),
    ("  .github/agents/qa-agent.md", 10, False, GRAY),
    ("  Role: multi-turn Q&A on the document", 10, False, DARK),
    ("  Tools: read_file + ask.py invocation", 10, False, GRAY),
    ("", 5, False, DARK),
    ("Bonus: Add summarizer agent", 11, True, NAVY),
])

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 17 — KEY TAKEAWAYS
# ─────────────────────────────────────────────────────────────────────────────
s = blank()
slide_header(s, "Key Takeaways")
footer(s)

kts = [
    ("Multi-Agent = Orchestrator + Specialists",
     "1 orchestrator, N specialist sub-agents, shared state.\n"
     "Each agent owns ONE domain. Start with 2, grow as needed."),
    ("GitHub Copilot Agents is a real framework",
     "Zero infra. Markdown-defined. MCP-compatible.\n"
     ".github/agents runs today inside VS Code — no setup."),
    ("Choose multi-agent when tasks decompose",
     "Context limits? Parallel subtasks? Composable tools?\n"
     "YES to any → multi-agent. Single context fits → single LLM."),
    ("RAG evolved to remove complexity layers",
     "V1 rich but heavy → V2 lean but keyword-only →\n"
     "V3 Agentic: zero retrieval code, LLM reads full doc."),
    ("Agentic RAG best for structured docs",
     "JSON tree + LLM reasoning = section-aware answers.\n"
     "Add .github/agents for multi-turn Q&A sessions."),
    ("Both use cases share the same principle",
     "Agents + tools + clear contracts = composable AI.\n"
     "GitHub Copilot Agents powers both patterns."),
]
KW3 = 5.88
for i, (title, body) in enumerate(kts):
    ci = i % 3;  ri = i // 3
    x = 0.35 + ci * 4.32;  y = 1.05 + ri * 3.00
    box(s, x, y, 4.10, 2.80, CARD)
    hline(s, x, y, 4.10, ORANGE, 0.06)
    box(s, x, y, 0.06, 2.80, ORANGE)
    tb(s, x+0.18, y+0.12, 3.84, 0.44, title, sz=12, bold=True, col=NAVY)
    hline(s, x+0.18, y+0.58, 3.74, LINE, 0.03)
    tb(s, x+0.18, y+0.68, 3.84, 1.96, body, sz=10, col=DARK)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 18 — THANK YOU
# ─────────────────────────────────────────────────────────────────────────────
s = blank()
box(s, 0, 0, 13.33, 7.50, WHITE)
box(s, 0, 0, 0.22, 7.50, ORANGE)
hline(s, 0.22, 3.55, 13.11, ORANGE, 0.06)

tb(s, 0.55, 0.75, 12.0, 1.35, "Thank You!", sz=46, bold=True, col=NAVY)
tb(s, 0.55, 2.20, 8.0, 0.46, "Questions & Discussion Welcome", sz=16, col=ORANGE)
hline(s, 0.55, 2.78, 5.5, LINE, 0.03)
tb(s, 0.55, 2.96, 10.0, 0.42,
   "Raviteja Thota  |  Senior Software Engineer  |  EPAM", sz=13, col=DARK)
tb(s, 0.55, 3.42, 10.0, 0.36,
   "EngX AI Coach — Final Exam  |  May 2026", sz=11, col=GRAY)

# Two use-case recap boxes
box(s, 0.55, 3.95, 5.70, 1.55, UC1_L)
hline(s, 0.55, 3.95, 5.70, UC1_C, 0.05)
tb(s, 0.72, 4.04, 5.38, 0.38, "Use Case 1  —  Multi-Agentic Solution", sz=11, bold=True, col=UC1_C)
tb(s, 0.72, 4.46, 5.38, 0.88,
   "Orchestrator + 5 agents  ·  Java toggle analysis\n.github/agents  ·  Confluence MCP  ·  MPGS",
   sz=10, col=DARK)

box(s, 0.55, 5.65, 5.70, 1.55, UC2_L)
hline(s, 0.55, 5.65, 5.70, UC2_C, 0.05)
tb(s, 0.72, 5.74, 5.38, 0.38, "Use Case 2  —  RAG Evolution Journey", sz=11, bold=True, col=UC2_C)
tb(s, 0.72, 6.16, 5.38, 0.88,
   "Vector RAG  ›  Vector-less  ›  Agentic RAG\nLLM as retriever  ·  3 Python files",
   sz=10, col=DARK)

tb(s, 0.55, 7.22, 12.0, 0.24,
   "OrgEngXSoftwareDevelopersCoach@epam.com  |  EPAM Proprietary & Confidential.",
   sz=9, col=LGRAY)

# ─────────────────────────────────────────────────────────────────────────────
prs.save(OUT)
print(f"Saved  ({len(prs.slides)} slides)  ->  {OUT}")
