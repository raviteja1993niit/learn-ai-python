"""
generate_ppt.py
Creates the EngX AI Coach Final Exam presentation from a blank python-pptx
Presentation — no template base, no placeholder collisions, no overlapping.

Slide size : 13.33 x 7.5 inches (widescreen 16:9)
Theme      : EPAM dark (navy background, orange accents, white text)
Run        : python generate_ppt.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Colour palette ────────────────────────────────────────────────────────────
BG      = RGBColor(0x0D, 0x0D, 0x2E)   # slide background  (dark navy)
ORANGE  = RGBColor(0xFA, 0x64, 0x00)   # EPAM orange
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY   = RGBColor(0xCC, 0xCC, 0xCC)
GRAY    = RGBColor(0x77, 0x77, 0x77)
GREEN   = RGBColor(0x4C, 0xAF, 0x50)
RED     = RGBColor(0xEF, 0x53, 0x50)
CODE_BG = RGBColor(0x14, 0x14, 0x28)
CODE_FG = RGBColor(0xCC, 0xDD, 0xEE)
CARD_BG = RGBColor(0x18, 0x18, 0x42)

OUT = (r"C:\Users\e135408\Downloads\personal-work\learn-ai\projects"
       r"\engx-ai-coach-for-software-developers"
       r"\presentations\engx-ai-coach-final-exam-raviteja-thota.pptx")

# ── Build presentation ────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.50)

SW = prs.slide_width
SH = prs.slide_height


# ─────────────────────────────────────────────────────────────────────────────
# Helper utilities
# ─────────────────────────────────────────────────────────────────────────────

def blank_slide():
    """Add a truly blank slide and paint the background."""
    layout = prs.slide_layouts[6]          # Blank layout
    slide  = prs.slides.add_slide(layout)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG
    return slide


def rect(slide, l, t, w, h, color):
    """Add a filled rectangle (no outline)."""
    shp = slide.shapes.add_shape(
        1,                                 # MSO_SHAPE_TYPE.RECTANGLE = 1
        Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp


def textbox(slide, l, t, w, h, text,
            size=14, bold=False, color=WHITE,
            align=PP_ALIGN.LEFT, wrap=True, italic=False):
    """Add a single-run text box."""
    txb = slide.shapes.add_textbox(
        Inches(l), Inches(t), Inches(w), Inches(h))
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text           = text
    run.font.size      = Pt(size)
    run.font.bold      = bold
    run.font.italic    = italic
    run.font.color.rgb = color
    return txb


def multiline(slide, l, t, w, h, rows,
              default_size=13, default_color=WHITE,
              default_bold=False, bg=None, pad=0.05):
    """
    Add a text box with multiple paragraphs.
    rows = list of str  OR  (text, size, bold, color)
    """
    txb = slide.shapes.add_textbox(
        Inches(l), Inches(t), Inches(w), Inches(h))
    if bg:
        txb.fill.solid()
        txb.fill.fore_color.rgb = bg
    tf = txb.text_frame
    tf.word_wrap = True
    tf.margin_left   = Inches(pad)
    tf.margin_right  = Inches(pad)
    tf.margin_top    = Inches(pad)
    tf.margin_bottom = Inches(pad)

    first = True
    for row in rows:
        if isinstance(row, str):
            txt, sz, bd, col = row, default_size, default_bold, default_color
        else:
            txt = row[0]
            sz  = row[1] if len(row) > 1 else default_size
            bd  = row[2] if len(row) > 2 else default_bold
            col = row[3] if len(row) > 3 else default_color
        para = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = para.add_run()
        run.text           = txt
        run.font.size      = Pt(sz)
        run.font.bold      = bd
        run.font.color.rgb = col
    return txb


def header_bar(slide, title, subtitle=""):
    """Paint the top header strip + title text."""
    rect(slide, 0, 0, 13.33, 1.05, BG)       # top strip (same as bg — clean separation)
    rect(slide, 0, 1.0,  13.33, 0.06, ORANGE) # thin orange rule
    textbox(slide, 0.35, 0.08, 12.5, 0.80,
            title, size=24, bold=True, color=WHITE)
    if subtitle:
        textbox(slide, 0.35, 0.72, 12.5, 0.30,
                subtitle, size=12, bold=False, color=ORANGE)


def footer(slide, text="EPAM Proprietary & Confidential."):
    rect(slide, 0, 7.25, 13.33, 0.25, RGBColor(0x08, 0x08, 0x1E))
    textbox(slide, 0.35, 7.27, 12.5, 0.20,
            text, size=9, color=GRAY)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — COVER
# ─────────────────────────────────────────────────────────────────────────────
s = blank_slide()
# big orange vertical accent bar on left
rect(s, 0, 0, 0.18, 7.50, ORANGE)
# title block
textbox(s, 0.55, 1.60, 12.0, 1.30,
        "How to Develop Your Own", size=36, bold=True, color=WHITE)
textbox(s, 0.55, 2.75, 12.0, 1.20,
        "Multi-Agentic Solution", size=36, bold=True, color=ORANGE)
# divider
rect(s, 0.55, 4.05, 6.0, 0.05, ORANGE)
# metadata
textbox(s, 0.55, 4.25, 9.0, 0.45,
        "Raviteja Thota  |  Senior Software Engineer  |  EPAM", size=14, color=WHITE)
textbox(s, 0.55, 4.72, 9.0, 0.38,
        "EngX AI Coach — Final Exam Presentation  |  May 2026", size=12, color=LGRAY)
# bottom tag
textbox(s, 0.55, 6.80, 9.0, 0.35,
        "EPAM Proprietary & Confidential.", size=10, color=GRAY)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — AGENDA
# ─────────────────────────────────────────────────────────────────────────────
s = blank_slide()
header_bar(s, "Agenda")
footer(s)
items = [
    "01   About Me",
    "02   Topic Introduction — What is a Multi-Agent System?",
    "03   AI Trend: The RAG Evolution  (Vector  ›  Vector-less  ›  Agentic)",
    "04   Theory — Orchestrator, Agents, Tools & Memory",
    "05   Framework Landscape  (LangChain / LangGraph / ADK / GitHub Copilot Agents)",
    "06   Real Use Case — Agentic Toggle Management System  (Live Code)",
    "07   Code Demo — Agentic RAG App  (vector-less-rag-agentic)",
    "08   Practical Task for You  (15 min)",
    "09   Key Takeaways & Q&A",
]
for i, item in enumerate(items):
    col = ORANGE if i == 0 else WHITE
    sz  = 15     if i == 0 else 14
    bd  = True   if i == 0 else False
    textbox(s, 0.55, 1.18 + i * 0.64, 12.2, 0.55,
            item, size=sz, bold=bd, color=col)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — ABOUT ME
# ─────────────────────────────────────────────────────────────────────────────
s = blank_slide()
header_bar(s, "About Me")
footer(s)
# Left column
multiline(s, 0.35, 1.15, 6.1, 5.85, [
    ("Raviteja Thota", 22, True, ORANGE),
    ("Senior Software Engineer  @  EPAM", 13, False, WHITE),
    ("Hyderabad, India", 12, False, LGRAY),
    ("", 5, False, WHITE),
    ("Skills:", 13, True, ORANGE),
    ("Angular  ·  Java / Spring Boot  ·  AWS  ·  Microservices", 12, False, WHITE),
    ("", 5, False, WHITE),
    ("Certifications:", 13, True, ORANGE),
    ("  Claude Certified Architect – Foundations  (Anthropic)", 12, False, WHITE),
    ("  Microsoft AZ-400 DevOps Engineer Expert", 12, False, WHITE),
    ("  AWS Certified Developer – Associate", 12, False, WHITE),
    ("  Google Professional Cloud Developer", 12, False, WHITE),
    ("  Oracle OCI Generative AI Professional", 12, False, WHITE),
    ("  AWS AI Practitioner", 12, False, WHITE),
    ("", 5, False, WHITE),
    ("India Innovation Enabler  (Feb 2026)", 12, True, ORANGE),
    ("  MPGS AI Run Adoption — Smart Mapper,", 11, False, WHITE),
    ("  ATF to Flow Framework featured at LEAP (global)", 11, False, LGRAY),
])
# Right column
rect(s, 6.75, 1.15, 0.04, 5.85, ORANGE)   # vertical divider
multiline(s, 7.0, 1.15, 5.95, 5.85, [
    ("Apps Built for This Exam:", 13, True, ORANGE),
    ("", 5, False, WHITE),
    ("agentic-ai-toggle-management", 12, True, WHITE),
    ("  Orchestrator + 5 agents, Confluence MCP,", 11, False, LGRAY),
    ("  Java toggle analysis + code refactoring", 11, False, LGRAY),
    ("", 4, False, WHITE),
    ("vector-rag-app-v1", 12, True, WHITE),
    ("  ChromaDB + sentence-transformers + BM25", 11, False, LGRAY),
    ("", 4, False, WHITE),
    ("vector-less-rag-app-v2", 12, True, WHITE),
    ("  BM25 + FastAPI, no vector DB", 11, False, LGRAY),
    ("", 4, False, WHITE),
    ("vector-less-rag-app", 12, True, WHITE),
    ("  BM25 + Streamlit, multi-query search", 11, False, LGRAY),
    ("", 4, False, WHITE),
    ("vector-less-rag-agentic", 12, True, WHITE),
    ("  LLM as retriever + .github/agents", 11, False, LGRAY),
    ("", 4, False, WHITE),
    ("transcript-notes-generator", 12, True, WHITE),
    ("  faster-whisper + Copilot LLM + agents", 11, False, LGRAY),
])

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — TOPIC INTRODUCTION
# ─────────────────────────────────────────────────────────────────────────────
s = blank_slide()
header_bar(s, "What is a Multi-Agent System?")
footer(s)
multiline(s, 0.35, 1.18, 12.6, 5.85, [
    ("A system where multiple AI agents collaborate to solve complex tasks.",
     16, True, ORANGE),
    ("Each agent has a specific role, toolset, and output contract.", 13, False, WHITE),
    ("An Orchestrator routes work.  Sub-agents execute specialised tasks.", 13, False, WHITE),
    ("", 5, False, WHITE),
    ("Key Principles:", 13, True, ORANGE),
    ("  Separation of Concerns  —  each agent does ONE thing well", 13, False, WHITE),
    ("  Composability           —  agents can be added or swapped independently", 13, False, WHITE),
    ("  Transparency            —  every delegation step is logged and traceable", 13, False, WHITE),
    ("  Guardrails              —  orchestrator enforces what each agent CANNOT do", 13, False, WHITE),
    ("", 5, False, WHITE),
    ("Why not just use a single LLM?", 13, True, ORANGE),
    ("  Context window limits  —  one LLM cannot see 16,000+ files at once", 13, False, RED),
    ("  No specialisation       —  a generalist LLM is worse than a focused agent", 13, False, RED),
    ("  Hard to maintain        —  monolithic prompts break as complexity grows", 13, False, RED),
    ("", 5, False, WHITE),
    ("Target audience today:  Developers using GenAI for 3 months, ready for the next level.",
     12, True, GREEN),
])

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — RAG EVOLUTION
# ─────────────────────────────────────────────────────────────────────────────
s = blank_slide()
header_bar(s, "AI Trend: The RAG Evolution",
           "Your 3 apps tell the story  —  from complex infra to zero-code retrieval")
footer(s)
CW, CH = 3.85, 5.35
starts = [0.35, 4.74, 9.13]
titles = ["V1 — Vector RAG", "V2 — Vector-less RAG", "V3 — Agentic RAG"]
apps   = ["vector-rag-app-v1", "vector-less-rag-app-v2", "vector-less-rag-agentic"]
bullets = [
    [("  ChromaDB embeddings",12,False,WHITE),
     ("  sentence-transformers",12,False,WHITE),
     ("  BM25 hybrid search",12,False,WHITE),
     ("",4,False,WHITE),
     ("+ Good semantic recall",11,False,GREEN),
     ("- Embedding setup cost",11,False,RED),
     ("- Vector DB infra needed",11,False,RED),
     ("- Complex codebase",11,False,RED)],
    [("  BM25 keyword search",12,False,WHITE),
     ("  No embeddings",12,False,WHITE),
     ("  Multi-query retrieval",12,False,WHITE),
     ("",4,False,WHITE),
     ("+ Zero infra, free",11,False,GREEN),
     ("+ Fast startup",11,False,GREEN),
     ("- Misses semantic matches",11,False,RED),
     ("- Still needs retriever code",11,False,RED)],
    [("  LLM IS the retriever",12,False,WHITE),
     ("  Full JSON doc as context",12,False,WHITE),
     ("  .github/agents for Q&A",12,False,WHITE),
     ("",4,False,WHITE),
     ("+ No retrieval code at all",11,False,GREEN),
     ("+ Structure-aware answers",11,False,GREEN),
     ("+ Only 3 Python files",11,False,GREEN),
     ("- Token limit on huge docs",11,False,RED)],
]
for i, x in enumerate(starts):
    rect(s, x, 1.18, CW, CH, CARD_BG)
    textbox(s, x+0.12, 1.28, CW-0.24, 0.50,
            titles[i], size=14, bold=True, color=ORANGE)
    textbox(s, x+0.12, 1.80, CW-0.24, 0.28,
            apps[i], size=10, bold=False, color=LGRAY)
    rect(s, x+0.12, 2.10, CW-0.24, 0.03, ORANGE)
    multiline(s, x+0.12, 2.20, CW-0.24, 4.0, bullets[i])
# arrows between cards
for ax in [4.28, 8.67]:
    textbox(s, ax, 3.55, 0.40, 0.40, ">", size=20, bold=True,
            color=ORANGE, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — THEORY: CORE CONCEPTS
# ─────────────────────────────────────────────────────────────────────────────
s = blank_slide()
header_bar(s, "Theory: Building Blocks of a Multi-Agent System")
footer(s)
concepts = [
    ("Orchestrator (Router)",
     "The brain. Receives user input, decides which sub-agent handles it. "
     "NEVER does leaf work — only routes, coordinates, and enforces guardrails."),
    ("Sub-Agents (Specialists)",
     "Each has one focused role: analyse, report, fetch data, advise, or manage state. "
     "Has its own tools, SKILL.md definitions, and a clear output contract."),
    ("Tools",
     "Functions agents can call: read_file, grep_search, run_in_terminal, "
     "mcp-confluence/getPage, insert_edit_into_file, run_subagent, ..."),
    ("Shared State / Memory",
     "Persistent storage across agent turns: CSV registries, JSON files, session workspace. "
     "Enables resumable, long-running workflows."),
    ("Skills  (SKILL.md)",
     "Reusable, documented capabilities inside an agent. Step-by-step instructions. "
     "e.g. java-source-search/SKILL.md  |  decomposition-planner/SKILL.md"),
]
CW2, CH2 = 2.30, 5.60
for i, (title, body) in enumerate(concepts):
    x = 0.35 + i * 2.55
    rect(s, x, 1.18, CW2, CH2, CARD_BG)
    textbox(s, x+0.10, 1.28, CW2-0.20, 0.45,
            title, size=12, bold=True, color=ORANGE)
    rect(s, x+0.10, 1.73, CW2-0.20, 0.03, ORANGE)
    multiline(s, x+0.10, 1.82, CW2-0.20, 4.70, [(body, 11, False, WHITE)])

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7 — FRAMEWORK LANDSCAPE
# ─────────────────────────────────────────────────────────────────────────────
s = blank_slide()
header_bar(s, "Framework Landscape: Which Multi-Agent Framework to Use?")
footer(s)
# Table header row
rect(s, 0.35, 1.18, 12.6, 0.40, ORANGE)
cols_x = [0.40, 2.55, 4.70, 8.60]
cols_w = [2.10, 2.10, 3.85, 4.05]
headers = ["Framework", "Language", "Strength", "Our Choice"]
for x, w, h in zip(cols_x, cols_w, headers):
    textbox(s, x, 1.20, w, 0.36, h, size=12, bold=True, color=WHITE)

rows_data = [
    ("LangChain",   "Python",   "Rich ecosystem, chains, memory",              "No"),
    ("LangGraph",   "Python",   "Graph workflows, cycles, checkpoints",        "No"),
    ("Google ADK",  "Python",   "Google-native, Gemini, MCP built-in",         "No  (plan future)"),
    ("CrewAI",      "Python",   "Role-based, easy YAML config",                "No"),
    ("AutoGen",     "Python",   "Microsoft, conversation-driven agents",       "No"),
    ("GH Copilot\nAgents", "Markdown", ".github/agents, native IDE, MCP-ready", "YES — all 6 apps"),
]
row_colors = [CARD_BG, BG, CARD_BG, BG, CARD_BG, RGBColor(0x1A, 0x2A, 0x10)]
for ri, (row, bg_c) in enumerate(zip(rows_data, row_colors)):
    y = 1.62 + ri * 0.58
    rect(s, 0.35, y, 12.6, 0.54, bg_c)
    for xi, (x, w, cell) in enumerate(zip(cols_x, cols_w, row)):
        col = ORANGE if ri == 5 else WHITE
        textbox(s, x, y+0.05, w, 0.46, cell, size=11, bold=(ri==5), color=col)

textbox(s, 0.35, 7.08, 12.5, 0.20,
        "GitHub Copilot Agents = zero extra infra, Markdown-defined, MCP-compatible, runs today in VS Code",
        size=10, bold=False, color=LGRAY)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8 — REAL USE CASE: PROBLEM
# ─────────────────────────────────────────────────────────────────────────────
s = blank_slide()
header_bar(s, "Real Use Case: Agentic Toggle Management System",
           "MPGS Payment Gateway — Mastercard")
footer(s)
# Problem box
rect(s, 0.35, 1.18, 12.6, 1.85, CARD_BG)
textbox(s, 0.50, 1.25, 2.0, 0.38,
        "The Problem", size=13, bold=True, color=ORANGE)
multiline(s, 0.50, 1.65, 12.30, 1.20, [
    ("100+ feature toggles across 7 Java repositories (~16,000 files total).", 12, False, WHITE),
    ('"What breaks if toggle X is ON vs OFF?" — nobody knew without reading Java for days.',
     12, False, WHITE),
    ("Confluence docs were stale. Manual analysis was error-prone and took days per toggle.",
     12, False, LGRAY),
])
# Solution header
textbox(s, 0.35, 3.15, 12.6, 0.38,
        "The Solution: 1 Orchestrator + 5 Sub-Agents", size=13, bold=True, color=ORANGE)
# Agent cards
agents = [
    ("toggle-\norchestrator", "Routes commands\nGuards scope\nAnswers inline"),
    ("toggle-analysis-\nengine  [SA-1]",  "Java source scan\nON/OFF impact\nS2A enrichment"),
    ("toggle-report-\nwriter  [SA-2]",    "6 report types\nMarkdown output\nCoverage checks"),
    ("toggle-confluence-\nenquiry  [SA-3]","4-tier MCP search\nPage metadata\nBatch fetch"),
    ("toggle-advisor\n[SA-4]  on-req",    "Decomp planning\nJava refactor\nToggle removal"),
]
AW = 2.30
for i, (name, desc) in enumerate(agents):
    ax = 0.35 + i * 2.55
    clr = ORANGE if i == 0 else CARD_BG
    rect(s, ax, 3.60, AW, 2.95, clr if i==0 else CARD_BG)
    rect(s, ax, 3.60, AW, 0.03, ORANGE)
    textbox(s, ax+0.08, 3.68, AW-0.16, 0.70,
            name, size=11, bold=True, color=WHITE if i else BG)
    rect(s, ax+0.08, 4.38, AW-0.16, 0.02, WHITE if i else BG)
    multiline(s, ax+0.08, 4.45, AW-0.16, 1.90,
              [(l, 11, False, WHITE if i else BG) for l in desc.split("\n")])

textbox(s, 0.35, 6.65, 12.6, 0.35,
        "Result: Full toggle impact analysis for all 100+ toggles in minutes, not days.",
        size=12, bold=True, color=GREEN)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 9 — ARCHITECTURE DIAGRAM
# ─────────────────────────────────────────────────────────────────────────────
s = blank_slide()
header_bar(s, "Architecture: Orchestrator + 5 Sub-Agents")
footer(s)
arch = """\
USER COMMAND  (plain English or /slash-command)
       |
       v
+--------------------------------+
|      TOGGLE-ORCHESTRATOR       |  <-- Session entry point
|  Routes  *  Answers  *  Guards |      NEVER does leaf work
+-----+--------+--------+--------+
      |        |        |        |
      v        v        v        v
  [SA-1]    [SA-2]   [SA-3]   [SA-5]
  ANALYSIS  REPORT   CONFLU-  WORKSPACE
  ENGINE    WRITER   ENCE     MANAGER
                     ENQUIRY
  Java scan  6 types  4-tier   Session
  ON/OFF     Markdown MCP      state &
  impact     output   search   registry

                    [SA-4] ADVISOR  (* on-request only)
                     Decomp planning
                     Java code refactor / toggle removal

Tools available per agent:
  read_file  |  grep_search  |  run_subagent  |  insert_edit_into_file
  mcp-confluence/getPage  |  run_in_terminal  |  validate_cves  |  ...
"""
rect(s, 0.35, 1.18, 12.6, 5.85, CODE_BG)
textbox(s, 0.55, 1.28, 12.20, 5.65, arch,
        size=11, color=CODE_FG, wrap=True)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 10 — CODE: AGENT DEFINITION
# ─────────────────────────────────────────────────────────────────────────────
s = blank_slide()
header_bar(s, "Code: Defining an Agent",
           ".github/agents/toggle-orchestrator.agent.md")
footer(s)
code1 = """\
---
name: toggle-orchestrator
description: Session entry point and router for the toggle management system.
tools:
  - read_file
  - run_subagent          # delegates to sub-agents
  - grep_search
  - mcp-confluence/getPage
  - insert_edit_into_file
---

# Toggle Orchestrator

## Routing Strategy

### Tier 1 — Handle Inline  (no sub-agent needed)
| Request                    | Action                                      |
|----------------------------|---------------------------------------------|
| Single toggle data query   | Read CSV, reply as Markdown                 |
| Analysis status check      | Count registry rows, answer inline          |
| Toggle already analysed    | Display from analysis CSV directly          |

### Tier 2 — Delegate to Sub-Agent
| User Command                         | Sub-Agent                      |
|--------------------------------------|--------------------------------|
| analyze toggle <name>                | toggle-analysis-engine  (SA-1) |
| generate report                      | toggle-report-writer    (SA-2) |
| fetch confluence data for <name>     | toggle-confluence-enquiry SA-3 |
| plan toggle decomposition for <repo> | toggle-advisor          (SA-4) |
"""
rect(s, 0.35, 1.18, 12.6, 5.85, CODE_BG)
textbox(s, 0.50, 1.28, 12.20, 5.65, code1, size=10, color=CODE_FG)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 11 — CODE DEMO: AGENTIC RAG
# ─────────────────────────────────────────────────────────────────────────────
s = blank_slide()
header_bar(s, "Code Demo: Agentic RAG — The LLM IS the Retriever",
           "vector-less-rag-agentic / ask.py")
footer(s)
code2 = """\
# ask.py  —  The LLM reads the full document and retrieves relevant info
# No BM25.  No embeddings.  No vector DB.  3 Python files total.

def ask(json_path: str, question: str, model: str = "gpt-4o-mini") -> str:
    doc     = json.loads(Path(json_path).read_text(encoding="utf-8"))
    context = build_context(doc)   # full JSON  -->  readable Markdown sections

    client = openai.OpenAI(
        base_url = "https://models.inference.ai.azure.com",
        api_key  = get_token_or_raise(),   # reads from  gh auth token  automatically
    )
    response = client.chat.completions.create(
        model    = model,
        messages = [
            {"role": "system", "content": SYSTEM_PROMPT},
            {"role": "user",   "content": f"{context}\\n\\n{question}"},
        ],
        temperature = 0.2,
        max_tokens  = 1024,
    )
    return response.choices[0].message.content


# ── .github/agents/qa-agent.md  orchestrates multi-turn Q&A sessions ─────────
#
#  Role:   Question-answering specialist
#  Tools:  read_file  (loads the JSON)  +  calls ask.py
#  Flow:   user question  -->  qa-agent  -->  ask.py  -->  LLM  -->  answer
#
# ─────────────────────────────────────────────────────────────────────────────
# Traditional RAG:  embed  ->  store in ChromaDB  ->  BM25 retrieval  ->  LLM
# Agentic RAG:      convert doc to JSON  ->  send full JSON  ->  LLM retrieves
"""
rect(s, 0.35, 1.18, 12.6, 5.85, CODE_BG)
textbox(s, 0.50, 1.28, 12.20, 5.65, code2, size=10, color=CODE_FG)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 12 — DECISION GUIDE
# ─────────────────────────────────────────────────────────────────────────────
s = blank_slide()
header_bar(s, "Decision Guide: Single LLM  vs  Multi-Agent System")
footer(s)
# Table header
rect(s, 0.35, 1.18, 12.6, 0.40, ORANGE)
th = ["Scenario", "Single LLM", "Multi-Agent"]
txs = [0.40, 5.30, 9.80]
tws = [4.85, 4.45, 3.45]
for x, w, h in zip(txs, tws, th):
    textbox(s, x, 1.20, w, 0.36, h, size=12, bold=True, color=WHITE)

rows12 = [
    ("Simple Q&A, chatbot",                    "YES",             "No — overkill"),
    ("Document summarisation",                 "YES",             "YES"),
    ("Code analysis — 1 file",                 "YES",             "No — overkill"),
    ("Cross-repo scan (1000s of files)",        "No — ctx limit",  "YES — delegate per repo"),
    ("Parallel specialised subtasks",           "No — sequential", "YES — parallel agents"),
    ("Long-running, resumable workflow",        "No — timeout",    "YES — state persists"),
    ("Composable MCP / tool sets per role",     "No — one toolset","YES — each agent own tools"),
]
for ri, row in enumerate(rows12):
    y = 1.62 + ri * 0.62
    bg = CARD_BG if ri % 2 == 0 else BG
    rect(s, 0.35, y, 12.6, 0.58, bg)
    for xi, (x, w, cell) in enumerate(zip(txs, tws, row)):
        col = ORANGE if xi == 2 and "YES" in cell else WHITE if "YES" in cell else LGRAY
        textbox(s, x, y+0.06, w, 0.48, cell, size=11, color=col)

textbox(s, 0.35, 6.68, 12.6, 0.34,
        "Rule of thumb:  Can your task decompose into specialists with clear outputs?  YES -> Multi-agent.  NO -> Single LLM.",
        size=11, bold=True, color=ORANGE)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 13 — PRACTICAL TASK
# ─────────────────────────────────────────────────────────────────────────────
s = blank_slide()
header_bar(s, "Practical Task  —  Build Your Own Agent  (15 min)")
footer(s)
textbox(s, 0.35, 1.18, 12.6, 0.45,
        "Task: Create a 2-agent system that reviews a Python file for code smells.",
        size=14, bold=True, color=ORANGE)

steps = [
    ("Step 1", "Create  .github/agents/code-review-orchestrator.agent.md",
     "Role: Accept a file path from the user. Delegate to code-analyser. Return a summary report.\n"
     "Tools: run_subagent"),
    ("Step 2", "Create  .github/agents/code-analyser.agent.md",
     "Role: Read the file. Find: functions > 30 lines, no error handling, magic numbers.\n"
     "Tools: read_file, grep_search\n"
     "Output: Markdown table — issue | line | severity | suggestion"),
    ("Step 3", "Test in GitHub Copilot Chat (Agent Mode)",
     "@code-review-orchestrator  analyse src/app.py"),
]
for i, (label, title, body) in enumerate(steps):
    y = 1.75 + i * 1.60
    rect(s, 0.35, y, 12.6, 1.48, CARD_BG)
    textbox(s, 0.48, y+0.08, 1.20, 0.40, label, size=12, bold=True, color=ORANGE)
    textbox(s, 1.72, y+0.08, 11.10, 0.40, title, size=12, bold=True, color=WHITE)
    textbox(s, 1.72, y+0.52, 11.10, 0.84, body, size=11, color=LGRAY)

rect(s, 0.35, 6.60, 12.6, 0.45, RGBColor(0x0A, 0x25, 0x0A))
textbox(s, 0.50, 6.65, 12.3, 0.38,
        "Bonus: Add a 3rd agent that generates a JIRA-style ticket for each issue found.",
        size=12, bold=True, color=GREEN)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 14 — KEY TAKEAWAYS
# ─────────────────────────────────────────────────────────────────────────────
s = blank_slide()
header_bar(s, "Key Takeaways")
footer(s)
kts = [
    ("RAG has evolved — Vector  >  Vector-less  >  Agentic",
     "Each generation is simpler to build and more powerful to use. Start Agentic."),
    ("Multi-Agent = Orchestrator + Specialists + Tools + Shared State",
     "Not a buzzword — a structured way to organise multiple focused LLM calls."),
    ("GitHub Copilot Agents is a real, production multi-agent framework",
     "Zero infra. Markdown-defined. MCP-compatible. Runs today inside VS Code."),
    ("Start with the problem — then choose the framework",
     "ADK / LangGraph / CrewAI are all valid. Pick what fits your stack and team."),
    ("Agents need clear roles, contracts, and guardrails",
     "Define what each agent CAN and CANNOT do. Test each agent in isolation."),
    ("You already have the skills — start building",
     "A 2-agent system today. Scale up as your problem complexity grows."),
]
KW = 5.85
for i, (title, body) in enumerate(kts):
    col = i % 3
    row = i // 3
    x = 0.35 + col * 6.32
    y = 1.18 + row * 2.85
    rect(s, x, y, KW, 2.65, CARD_BG)
    rect(s, x, y, 0.06, 2.65, ORANGE)
    textbox(s, x+0.20, y+0.12, KW-0.30, 0.50,
            title, size=12, bold=True, color=ORANGE)
    rect(s, x+0.20, y+0.62, KW-0.30, 0.03, GRAY)
    multiline(s, x+0.20, y+0.72, KW-0.30, 1.75, [(body, 11, False, WHITE)])

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 15 — THANK YOU
# ─────────────────────────────────────────────────────────────────────────────
s = blank_slide()
rect(s, 0, 0, 0.18, 7.50, ORANGE)           # left accent bar
rect(s, 0.18, 3.55, 13.15, 0.05, ORANGE)    # mid horizontal rule
textbox(s, 0.55, 1.40, 12.0, 1.50,
        "Thank You!", size=48, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
textbox(s, 0.55, 3.00, 12.0, 0.55,
        "Questions & Discussion Welcome", size=16, bold=False, color=ORANGE,
        align=PP_ALIGN.LEFT)
textbox(s, 0.55, 3.80, 5.5, 0.40,
        "Raviteja Thota", size=15, bold=True, color=WHITE)
textbox(s, 0.55, 4.22, 8.0, 0.35,
        "Senior Software Engineer  |  EPAM  |  Hyderabad, India", size=12, color=LGRAY)
textbox(s, 0.55, 4.75, 9.0, 0.35,
        "Exam topic: How to Develop Your Own Multi-Agentic Solution", size=12, color=LGRAY)
textbox(s, 0.55, 5.18, 10.0, 0.35,
        "OrgEngXSoftwareDevelopersCoach@epam.com", size=12, color=ORANGE)
textbox(s, 0.55, 6.90, 12.0, 0.30,
        "EPAM Proprietary & Confidential.", size=9, color=GRAY)

# ─────────────────────────────────────────────────────────────────────────────
# SAVE
# ─────────────────────────────────────────────────────────────────────────────
prs.save(OUT)
print(f"Saved  ({len(prs.slides)} slides)  ->  {OUT}")
